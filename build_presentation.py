"""
build_presentation.py
─────────────────────
Builds the full cybersecurity presentation using python-pptx.
14 slides covering: title, problem, pipeline, imbalance, SMOTE,
all 5 models, feature importance, drift monitoring, API demo, conclusion.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
import copy
import lxml.etree as etree
import os

# ── Paths ────────────────────────────────────────────────────────────────────
BASE   = r"C:\Users\jadal\Desktop\cybersecurity"
RPTS   = os.path.join(BASE, "reports")
DOCS   = os.path.join(BASE, "docs")
OUT    = os.path.join(DOCS, "presentation.pptx")
os.makedirs(DOCS, exist_ok=True)

def rpt(name): return os.path.join(RPTS, name)

# ── Colors ───────────────────────────────────────────────────────────────────
BG      = RGBColor(0x0A, 0x16, 0x28)
CARD    = RGBColor(0x16, 0x22, 0x36)
CARD2   = RGBColor(0x1A, 0x2D, 0x44)
GREEN   = RGBColor(0x00, 0xD6, 0x8F)
RED     = RGBColor(0xFF, 0x47, 0x57)
BLUE    = RGBColor(0x1A, 0x56, 0xDB)
ORANGE  = RGBColor(0xFF, 0xA5, 0x00)
PURPLE  = RGBColor(0x9B, 0x59, 0xB6)
CYAN    = RGBColor(0x34, 0x98, 0xDB)
YELLOW  = RGBColor(0xF1, 0xC4, 0x0F)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
MUTED   = RGBColor(0xC5, 0xD0, 0xDC)
DIM     = RGBColor(0x6B, 0x7E, 0x94)

# ── Slide size 16:9 ──────────────────────────────────────────────────────────
W = Inches(10)
H = Inches(5.625)

# ═════════════════════════════════════════════════════════════════════════════
#  LOW-LEVEL HELPERS
# ═════════════════════════════════════════════════════════════════════════════

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)

def bg(slide, color=BG):
    """Fill slide background."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, x, y, w, h, fill_color, line_color=None, line_width=Pt(0.5)):
    """Add a filled rectangle."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def oval(slide, x, y, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        9,  # MSO_SHAPE_TYPE.OVAL
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def line_shape(slide, x, y, w, h, color, width=Pt(1.5)):
    connector = slide.shapes.add_connector(
        1,  # STRAIGHT
        Inches(x), Inches(y), Inches(x + w), Inches(y + h)
    )
    connector.line.color.rgb = color
    connector.line.width = width
    return connector

def txb(slide, text, x, y, w, h,
        size=12, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
        italic=False, font="Calibri", wrap=True):
    """Add a text box."""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return txBox

def txb_para(slide, runs, x, y, w, h, size=12, align=PP_ALIGN.LEFT,
             font="Calibri", wrap=True):
    """Add text box with multiple runs (each run: (text, color, bold, italic))."""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    for (txt, col, bld, itl) in runs:
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(size)
        r.font.bold = bld
        r.font.italic = itl
        r.font.color.rgb = col
        r.font.name = font
    return txBox

def txb_bullets(slide, items, x, y, w, h, size=11, color=MUTED, font="Calibri"):
    """Add a text box with bullet points (list of strings)."""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        # bullet via XML
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        buChar = etree.SubElement(pPr, qn('a:buChar'))
        buChar.set('char', '▸')
        run = p.add_run()
        run.text = item
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = font
    return txBox

def add_image(slide, path, x, y, w, h):
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))

def add_title(slide, text):
    txb(slide, text, 0.4, 0.18, 9.2, 0.6, size=28, bold=True, color=WHITE)

def add_subtitle(slide, text):
    txb(slide, text, 0.4, 0.82, 9.2, 0.32, size=11, color=MUTED)

def divider(slide):
    line_shape(slide, 0.4, 0.78, 9.2, 0, GREEN, Pt(1.5))

def card_rect(slide, x, y, w, h, color=CARD):
    rect(slide, x, y, w, h, color, line_color=None)

def accent_bar(slide, x, y, h, color=GREEN):
    rect(slide, x, y, 0.06, h, color)

def stat_box(slide, x, y, value, label, color=GREEN):
    card_rect(slide, x, y, 2.2, 1.1)
    txb(slide, value, x, y + 0.05, 2.2, 0.62,
        size=34, bold=True, color=color, align=PP_ALIGN.CENTER)
    txb(slide, label, x, y + 0.68, 2.2, 0.35,
        size=10, color=MUTED, align=PP_ALIGN.CENTER)

def tag_box(slide, text, x, y, color=BLUE):
    rect(slide, x, y, 1.65, 0.3, CARD2, line_color=color)
    txb(slide, text, x, y, 1.65, 0.3, size=9, color=MUTED, align=PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════════════════════════
#  SLIDES
# ═════════════════════════════════════════════════════════════════════════════

def slide01(prs):
    """Title slide"""
    s = blank_slide(prs)
    bg(s)
    # Top bar
    rect(s, 0, 0, 10, 0.08, GREEN)

    txb(s, "IoT Intrusion Detection",
        0.5, 0.5, 9, 0.85, size=46, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb(s, "with Real-Time Drift Monitoring",
        0.5, 1.35, 9, 0.55, size=26, bold=False, color=GREEN, align=PP_ALIGN.CENTER)
    txb(s, "Multi-class ML system  ·  9 attack families  ·  CICIoT2023 dataset",
        0.5, 1.92, 9, 0.35, size=13, color=MUTED, align=PP_ALIGN.CENTER)

    stat_box(s, 0.5,  2.5,  "9",      "Attack Classes",  GREEN)
    stat_box(s, 2.85, 2.5,  "7.8M+",  "Raw Flows",       BLUE)
    stat_box(s, 5.2,  2.5,  "5",      "Models Trained",  ORANGE)
    stat_box(s, 7.55, 2.5,  "0.830",  "Best Macro-F1",   CYAN)

    tags = ["FastAPI", "Docker", "LightGBM ★", "PyTorch MLP", "KS Drift Test"]
    for i, t in enumerate(tags):
        tag_box(s, t, 0.5 + i * 1.82, 3.85, BLUE)

    # Bottom bar
    rect(s, 0, 5.3, 10, 0.32, RGBColor(0x0D, 0x1E, 0x35))
    txb(s, "CICIoT2023 · University of New Brunswick  |  Macro-F1 as primary metric (class imbalance)",
        0.4, 5.32, 9.2, 0.26, size=9, color=DIM, align=PP_ALIGN.CENTER)


def slide02(prs):
    """Problem statement"""
    s = blank_slide(prs)
    bg(s)
    add_title(s, "The Problem — IoT Networks Under Attack")
    divider(s)
    add_subtitle(s, "33 raw attack types → 9 families  ·  7.8M+ labeled flows  ·  extreme class imbalance")

    # Left card
    card_rect(s, 0.4, 1.05, 4.4, 3.9)
    accent_bar(s, 0.4, 1.05, 3.9, RED)
    txb(s, "Why this matters", 0.6, 1.12, 4.0, 0.38, size=15, bold=True, color=RED)
    problems = [
        "Smart TVs, cameras, sensors — billions of devices with zero built-in security",
        "Attackers exploit them for DDoS, botnet C2, and lateral movement",
        "Traditional signature-based IDS fails on novel attack patterns",
        "Class imbalance: 27.9x more Benign than rare attack types",
        "We need ML that generalises across 9 different threat families",
    ]
    txb_bullets(s, problems, 0.6, 1.55, 4.0, 3.1, size=11, color=MUTED)

    # Right card
    card_rect(s, 5.05, 1.05, 4.5, 3.9)
    accent_bar(s, 5.05, 1.05, 3.9, GREEN)
    txb(s, "9 Attack Families", 5.25, 1.12, 4.1, 0.38, size=15, bold=True, color=GREEN)

    classes = [
        ("Benign",               "Normal IoT traffic",              GREEN),
        ("DDoS",                 "Distributed volumetric floods",   RED),
        ("DoS",                  "Single-source service denial",    ORANGE),
        ("Mirai",                "Botnet C2 & propagation",         BLUE),
        ("Spoofing",             "DNS / ARP poisoning",             PURPLE),
        ("Recon",                "Port scans, OS fingerprinting",   ORANGE),
        ("Recon-HostDiscovery",  "Host discovery sweeps",           CYAN),
        ("Web",                  "SQLi, XSS, command injection",    CYAN),
        ("DictionaryBruteForce", "Password brute force",            RED),
    ]
    for i, (name, desc, col) in enumerate(classes):
        yy = 1.58 + i * 0.36
        rect(s, 5.25, yy, 0.08, 0.24, col)
        txb(s, name, 5.42, yy, 1.6, 0.26, size=10, bold=True, color=WHITE)
        txb(s, desc, 7.1,  yy, 2.3, 0.26, size=9,  color=MUTED)


def slide03(prs):
    """Dataset & pipeline"""
    s = blank_slide(prs)
    bg(s)
    add_title(s, "Dataset & ML Pipeline")
    divider(s)
    add_subtitle(s, "CICIoT2023  ·  46 features  ·  Stratified sampling → SMOTE → 5 models")

    steps = [
        ("1", "Raw CSV",   "7.8M flows\n33 labels",    BLUE),
        ("2", "Map Labels","33 → 9\nfamilies",          CYAN),
        ("3", "Sample",    "≤20K/class\nstratified",    ORANGE),
        ("4", "Split",     "70/15/15\ntemporal",        PURPLE),
        ("5", "SMOTE",     "Balance rare\nclasses",     GREEN),
        ("6", "Scale",     "StandardScaler\n+impute",   YELLOW),
        ("7", "Train",     "5 models\nColab GPU",       RED),
        ("8", "Deploy",    "FastAPI\n+Docker",          GREEN),
    ]
    for i, (num, label, detail, col) in enumerate(steps):
        x = 0.35 + i * 1.17
        oval(s, x + 0.22, 1.08, 0.65, 0.65, col, col)
        txb(s, num,    x + 0.22, 1.08, 0.65, 0.65, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txb(s, label,  x,        1.78, 1.1,  0.3,  size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txb(s, detail, x,        2.08, 1.1,  0.45, size=8,  color=MUTED, align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            line_shape(s, x + 1.0, 1.4, 0.18, 0, MUTED, Pt(1.2))

    decisions = [
        ("Why Macro-F1?",
         "Accuracy is deceptive with imbalanced data. A model predicting only 'Benign' scores 17% accuracy but catches zero attacks. Macro-F1 penalises equally across all 9 classes.",
         ORANGE),
        ("Why Temporal Split?",
         "Random shuffle leaks future data into training. We split by row order (time proxy) to prevent data leakage — the model never sees 'future' traffic during training.",
         BLUE),
        ("Why SMOTE?",
         "DictionaryBruteForce had only 1,541 training samples. SMOTE synthesises new minority examples by interpolating nearest neighbours, not simple copying.",
         GREEN),
    ]
    for i, (title, body, col) in enumerate(decisions):
        x = 0.4 + i * 3.1
        card_rect(s, x, 2.72, 2.9, 2.2)
        accent_bar(s, x, 2.72, 2.2, col)
        txb(s, title, x + 0.18, 2.78, 2.6, 0.36, size=13, bold=True, color=col)
        txb(s, body,  x + 0.18, 3.18, 2.6, 1.6,  size=9.5, color=MUTED)


def slide04(prs):
    """Class imbalance plot"""
    s = blank_slide(prs)
    bg(s)
    add_title(s, "Class Imbalance — The Core Challenge")
    divider(s)
    add_subtitle(s, "Raw dataset: 27.9x imbalance between largest and smallest class")

    add_image(s, rpt("class_distribution.png"), 0.4, 1.05, 9.2, 3.4)

    card_rect(s, 0.4, 4.55, 9.2, 0.72)
    txb_para(s, [
        ("Problem: ",      RED,   True,  False),
        ("Standard ML trains on majority class — minority attacks get ignored. A naive model always guessing 'Benign' scores 83% accuracy but catches ZERO attacks.  ",
                           MUTED, False, False),
        ("Fix: SMOTE + class_weight='balanced'",
                           GREEN, True,  False),
    ], 0.6, 4.6, 8.8, 0.62, size=10.5)


def slide05(prs):
    """SMOTE"""
    s = blank_slide(prs)
    bg(s)
    add_title(s, "Fixing Imbalance — SMOTE Oversampling")
    divider(s)
    add_subtitle(s, "Synthetic Minority Over-sampling: synthesises new examples, not simple copies")

    add_image(s, rpt("smote_comparison.png"), 0.4, 1.05, 9.2, 3.2)

    cards = [
        ("Step 1 — Pick minority sample", "Select a random sample from the minority class (e.g. DictionaryBruteForce)", ORANGE),
        ("Step 2 — Find K neighbours",    "Find its K nearest neighbours in feature space using Euclidean distance",    BLUE),
        ("Step 3 — Interpolate",          "Synthesise new sample by interpolating between the point and a neighbour",   GREEN),
    ]
    for i, (title, body, col) in enumerate(cards):
        x = 0.4 + i * 3.1
        card_rect(s, x, 4.35, 2.9, 1.0)
        accent_bar(s, x, 4.35, 1.0, col)
        txb(s, title, x + 0.18, 4.4,  2.6, 0.3,  size=11, bold=True, color=col)
        txb(s, body,  x + 0.18, 4.73, 2.6, 0.55, size=9.5, color=MUTED)


def slide06(prs):
    """Model 1 — Logistic Regression"""
    s = blank_slide(prs)
    bg(s)
    add_title(s, "Model 1 — Logistic Regression (Baseline)")
    divider(s)
    add_subtitle(s, "Softmax classifier  ·  L2 regularisation  ·  Isotonic calibration  ·  SAGA solver")

    add_image(s, rpt("confusion_matrix_logistic_regression.png"), 0.35, 1.05, 5.2, 3.85)

    stat_box(s, 5.85, 1.05, "0.484", "Macro-F1",     RED)
    stat_box(s, 7.75, 1.05, "0.929", "ROC-AUC",      ORANGE)
    stat_box(s, 5.85, 2.3,  "34.2%", "Benign FP%",   RED)
    stat_box(s, 7.75, 2.3,  "Worst", "vs All Models", DIM)

    card_rect(s, 5.85, 3.55, 3.7, 1.35, CARD2)
    accent_bar(s, 5.85, 3.55, 1.35, RED)
    txb(s, "Why LR struggled", 6.05, 3.6, 3.3, 0.33, size=13, bold=True, color=RED)
    txb_bullets(s, [
        "Linear boundary — IoT attacks are non-linear",
        "Recon & Web attacks confused with Benign",
        "34% of benign traffic falsely flagged as attack",
    ], 6.05, 3.95, 3.3, 0.85, size=10, color=MUTED)

    rect(s, 0.35, 4.95, 1.55, 0.28, RGBColor(0xFF, 0x47, 0x57))
    txb(s, "BASELINE", 0.35, 4.95, 1.55, 0.28, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def slide07(prs):
    """Model 2 — Random Forest"""
    s = blank_slide(prs)
    bg(s)
    add_title(s, "Model 2 — Random Forest")
    divider(s)
    add_subtitle(s, "200 trees  ·  Bagging + feature subsampling  ·  balanced_subsample  ·  OOB validation")

    add_image(s, rpt("confusion_matrix_random_forest.png"), 0.35, 1.05, 5.2, 3.85)

    stat_box(s, 5.85, 1.05, "0.808", "Macro-F1",      BLUE)
    stat_box(s, 7.75, 1.05, "0.994", "ROC-AUC",       GREEN)
    stat_box(s, 5.85, 2.3,  "9.5%",  "Benign FP%",    ORANGE)
    stat_box(s, 7.75, 2.3,  "+67%",  "vs LR Macro-F1", GREEN)

    card_rect(s, 5.85, 3.55, 3.7, 1.35, CARD2)
    accent_bar(s, 5.85, 3.55, 1.35, BLUE)
    txb(s, "Why RF leaps ahead", 6.05, 3.6, 3.3, 0.33, size=13, bold=True, color=BLUE)
    txb_bullets(s, [
        "200 trees vote — no single tree overfits",
        "balanced_subsample handles imbalance per tree",
        "OOB score used as held-out validation",
    ], 6.05, 3.95, 3.3, 0.85, size=10, color=MUTED)

    rect(s, 0.35, 4.95, 1.55, 0.28, BLUE)
    txb(s, "ENSEMBLE", 0.35, 4.95, 1.55, 0.28, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def slide08(prs):
    """Feature Importance"""
    s = blank_slide(prs)
    bg(s)
    add_title(s, "What the Model Looks At — Feature Importance")
    divider(s)
    add_subtitle(s, "Permutation importance on Random Forest  ·  Macro-F1 drop when feature is shuffled")

    add_image(s, rpt("feature_importance_random_forest.png"), 0.35, 1.05, 5.9, 4.25)

    callouts = [
        ("syn_flag_number", "SYN floods are the signature of DDoS/DoS — top discriminator",   RED),
        ("Rate",            "Packet rate immediately reveals volumetric attacks",              ORANGE),
        ("ack_flag_number", "ACK ratio separates normal TCP from attack traffic",              BLUE),
        ("flow_duration",   "Short flows = scans; Long flows = legitimate sessions",           GREEN),
        ("Header_Length",   "Crafted headers expose spoofing & injection attacks",             PURPLE),
    ]
    for i, (feat, desc, col) in enumerate(callouts):
        y = 1.15 + i * 0.82
        card_rect(s, 6.6, y, 3.0, 0.72)
        rect(s, 6.6, y, 0.07, 0.72, col)
        txb(s, feat, 6.78, y + 0.05, 2.7, 0.28, size=12, bold=True, color=col)
        txb(s, desc, 6.78, y + 0.35, 2.7, 0.3,  size=9,  color=MUTED)


def slide09(prs):
    """XGBoost vs LightGBM"""
    s = blank_slide(prs)
    bg(s)
    add_title(s, "Models 3 & 4 — XGBoost vs LightGBM")
    divider(s)
    add_subtitle(s, "Both are gradient boosting ensembles — LightGBM wins with leaf-wise growth")

    add_image(s, rpt("confusion_matrix_xgboost.png"),  0.3,  1.05, 4.55, 3.4)
    add_image(s, rpt("confusion_matrix_lightgbm.png"), 5.15, 1.05, 4.55, 3.4)

    rect(s, 0.3,  1.05, 2.3, 0.3, ORANGE)
    txb(s, "XGBoost  F1=0.815",      0.3,  1.05, 2.3, 0.3, size=10, bold=True, color=BG, align=PP_ALIGN.CENTER)

    rect(s, 5.15, 1.05, 2.55, 0.3, GREEN)
    txb(s, "LightGBM \u2605  F1=0.830", 5.15, 1.05, 2.55, 0.3, size=10, bold=True, color=BG, align=PP_ALIGN.CENTER)

    cmp = [
        ("XGBoost",   "Level-wise growth, 2nd-order gradients, exact split search. 300 estimators, lr=0.05, multi:softprob objective.", ORANGE),
        ("LightGBM \u2605", "Leaf-wise growth (deeper branches), GOSS sampling, histogram bins. 10x faster training. Best Macro-F1=0.830.", GREEN),
    ]
    for i, (title, body, col) in enumerate(cmp):
        x = 0.3 + i * 4.85
        card_rect(s, x, 4.52, 4.55, 0.82)
        accent_bar(s, x, 4.52, 0.82, col)
        txb(s, title, x + 0.18, 4.56, 4.1, 0.3, size=13, bold=True, color=col)
        txb(s, body,  x + 0.18, 4.87, 4.1, 0.42, size=9.5, color=MUTED)


def slide10(prs):
    """MLP Neural Network"""
    s = blank_slide(prs)
    bg(s)
    add_title(s, "Model 5 — MLP Neural Network (PyTorch)")
    divider(s)
    add_subtitle(s, "4-layer network  ·  BatchNorm  ·  Dropout  ·  Adam optimizer  ·  Early stopping on Colab T4 GPU")

    add_image(s, rpt("confusion_matrix_mlp.png"), 0.3, 1.05, 4.7, 3.5)

    # Architecture layers
    layers = [
        ("Input",     "46 features",                    MUTED,  1.15),
        ("Dense 256", "+ BatchNorm + ReLU + Dropout(0.3)", BLUE, 1.15),
        ("Dense 128", "+ BatchNorm + ReLU + Dropout(0.3)", BLUE, 1.0),
        ("Dense 64",  "+ BatchNorm + ReLU + Dropout(0.2)", BLUE, 0.85),
        ("Output",    "9 classes  (Softmax)",            GREEN,  0.7),
    ]
    for i, (label, desc, col, w) in enumerate(layers):
        y = 1.12 + i * 0.68
        rect(s, 5.25, y, w, 0.42, col)
        txb(s, label, 5.25, y, w, 0.42, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txb(s, desc,  5.25 + w + 0.1, y + 0.07, 3.5, 0.28, size=9, color=MUTED)
        if i < len(layers) - 1:
            line_shape(s, 5.25 + w/2, y + 0.42, 0, 0.26, DIM, Pt(1))

    stat_box(s, 5.25, 3.3, "0.630", "Macro-F1",  PURPLE)
    stat_box(s, 7.45, 3.3, "0.976", "ROC-AUC",   BLUE)

    card_rect(s, 5.25, 4.58, 4.35, 0.72, CARD2)
    accent_bar(s, 5.25, 4.58, 0.72, PURPLE)
    txb(s, "Surprise finding", 5.43, 4.62, 3.95, 0.28, size=12, bold=True, color=PURPLE)
    txb(s, "MLP underperforms LightGBM by 20% on Macro-F1 — tabular data favours gradient-boosted trees over deep learning",
        5.43, 4.9, 3.95, 0.35, size=9.5, color=MUTED)


def slide11(prs):
    """Model Comparison"""
    s = blank_slide(prs)
    bg(s)
    add_title(s, "Model Comparison — All 5 Models")
    divider(s)
    add_subtitle(s, "Test set: 102,167 samples  ·  LightGBM selected as production model")

    add_image(s, rpt("model_comparison_chart.png"), 0.3, 1.05, 9.4, 3.55)

    takeaways = [
        ("LR: linear boundary can't separate IoT attacks",         RED),
        ("RF / XGB / LGBM: tree ensembles dominate tabular data",  GREEN),
        ("MLP: needs more data / tuning for tabular tasks",        ORANGE),
    ]
    for i, (text, col) in enumerate(takeaways):
        x = 0.3 + i * 3.17
        rect(s, x, 4.7, 3.05, 0.55, CARD, line_color=col)
        txb(s, text, x + 0.1, 4.72, 2.85, 0.51, size=9.5, color=col)


def slide12(prs):
    """Drift Monitoring"""
    s = blank_slide(prs)
    bg(s)
    add_title(s, "Drift Monitoring — Is the Model Still Valid?")
    divider(s)
    add_subtitle(s, "3 complementary checks: KS test  ·  PSI  ·  3-sigma alert rate monitor")

    monitors = [
        ("KS Drift Test", "KS",
         "Kolmogorov-Smirnov test compares each feature's distribution between training data and new incoming traffic.\n\np-value < 0.05 means the feature has significantly shifted.\nRuns every 500 predictions.",
         "p < 0.05 = DRIFTED", BLUE),
        ("PSI Score", "PSI",
         "Population Stability Index measures distribution shift.\n\nPSI < 0.1   Stable\nPSI 0.1-0.2  Moderate change\nPSI > 0.2   Significant DRIFT",
         "PSI > 0.2 = ACTION", ORANGE),
        ("Alert Rate Monitor", "3\u03c3",
         "Tracks % of traffic flagged as attacks in rolling windows.\n\nSudden spike = real attack OR model has drifted.\n\nUses 3-sigma rule: Z-score > 3 triggers CRITICAL warning.",
         "Z > 3 = CRITICAL", RED),
    ]
    for i, (title, icon, body, thresh, col) in enumerate(monitors):
        x = 0.4 + i * 3.1
        card_rect(s, x, 1.05, 2.9, 3.5)
        accent_bar(s, x, 1.05, 3.5, col)
        oval(s, x + 1.0, 1.15, 0.85, 0.85, col)
        txb(s, icon,   x + 1.0, 1.15, 0.85, 0.85, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Consolas")
        txb(s, title,  x + 0.18, 2.12, 2.58, 0.36, size=13, bold=True, color=col)
        txb(s, body,   x + 0.18, 2.52, 2.58, 1.5,  size=9.5, color=MUTED)
        rect(s, x + 0.18, 4.1, 2.55, 0.28, col)
        txb(s, thresh, x + 0.18, 4.1, 2.55, 0.28, size=9.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    card_rect(s, 0.4, 4.55, 9.2, 0.68)
    txb_para(s, [
        ("API Integration: ", GREEN, True, False),
        ("Every /predict call feeds features into DriftMonitor.update()  →  drift checked every 500 calls  →  GET /drift/status returns full PSI + KS report", MUTED, False, False),
    ], 0.6, 4.6, 8.8, 0.58, size=10.5)


def slide13(prs):
    """FastAPI Demo"""
    s = blank_slide(prs)
    bg(s)
    add_title(s, "Live System — FastAPI + Docker Deployment")
    divider(s)
    add_subtitle(s, "3-step startup  ·  Swagger UI  ·  5 REST endpoints  ·  Docker Hub image")

    # Left: endpoints
    card_rect(s, 0.4, 1.05, 4.4, 3.25)
    accent_bar(s, 0.4, 1.05, 3.25, GREEN)
    txb(s, "REST Endpoints", 0.6, 1.1, 4.0, 0.36, size=14, bold=True, color=GREEN)

    endpoints = [
        ("POST", "/predict",        "Classify single flow → label + confidence", BLUE),
        ("POST", "/predict/batch",  "Classify multiple flows at once",           BLUE),
        ("GET",  "/drift/status",   "PSI + KS drift report",                     GREEN),
        ("GET",  "/health",         "Liveness check",                            GREEN),
        ("GET",  "/docs",           "Interactive Swagger UI",                    GREEN),
    ]
    for i, (method, ep, desc, col) in enumerate(endpoints):
        y = 1.56 + i * 0.54
        rect(s, 0.6, y, 0.65, 0.26, col)
        txb(s, method, 0.6, y, 0.65, 0.26, size=9, bold=True, color=BG, align=PP_ALIGN.CENTER)
        txb(s, ep,   1.32, y,        1.55, 0.26, size=10, bold=True, color=WHITE, font="Consolas")
        txb(s, desc, 1.32, y + 0.26, 3.3,  0.22, size=8.5, color=DIM)

    # Right: 3 steps
    card_rect(s, 5.05, 1.05, 4.5, 3.25)
    accent_bar(s, 5.05, 1.05, 3.25, BLUE)
    txb(s, "Start in 3 Steps", 5.25, 1.1, 4.1, 0.36, size=14, bold=True, color=BLUE)

    steps = [
        ("1", "git clone  github.com/…/cybersecurity",              MUTED),
        ("2", "pip install -r requirements.txt",                     MUTED),
        ("3", "py -3.10 -m uvicorn src.api:app --port 8000",         GREEN),
    ]
    for i, (n, cmd, col) in enumerate(steps):
        y = 1.56 + i * 0.52
        oval(s, 5.25, y, 0.36, 0.36, BLUE)
        txb(s, n, 5.25, y, 0.36, 0.36, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txb(s, cmd, 5.7, y + 0.04, 3.7, 0.28, size=9, color=col, font="Consolas")

    line_shape(s, 5.25, 3.18, 4.2, 0, DIM, Pt(0.5))
    txb(s, "Or with Docker:", 5.25, 3.22, 4.1, 0.26, size=10, bold=True, color=ORANGE)
    txb(s, "docker pull <username>/iot-ids:latest\ndocker run -p 8000:8000 <username>/iot-ids:latest",
        5.25, 3.5, 4.2, 0.45, size=9, color=MUTED, font="Consolas")

    # Sample response
    card_rect(s, 0.4, 4.38, 9.2, 0.87)
    txb_para(s, [
        ("Example response: ", GREEN, True, False),
        ('{ "label": "DDoS", "confidence": 0.9997, "model_used": "LightGBM", "probabilities": {"DDoS": 0.9997, ...} }',
         MUTED, False, False),
    ], 0.6, 4.44, 8.8, 0.75, size=10, font="Calibri")


def slide14(prs):
    """Conclusion"""
    s = blank_slide(prs)
    bg(s)
    rect(s, 0, 0, 10, 0.08, GREEN)

    txb(s, "Key Takeaways", 0.5, 0.25, 9, 0.62,
        size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    takeaways = [
        ("01", "Trees beat deep learning on tabular data",
         "LightGBM outperforms a 4-layer PyTorch MLP by 20% Macro-F1. Tabular IoT features don't benefit from raw neural feature learning — gradient boosting wins.", GREEN),
        ("02", "Imbalance must be fixed, not ignored",
         "SMOTE + class_weight='balanced' prevents the model collapsing to majority class. Without it, rare attacks like DictionaryBruteForce are invisible.", ORANGE),
        ("03", "Accuracy is a misleading metric",
         "A model predicting 'Benign' every time scores 83% accuracy but catches zero attacks. Macro-F1 is the honest metric for imbalanced multi-class problems.", RED),
        ("04", "Drift monitoring is production-critical",
         "KS test + PSI + 3-sigma alert rate catches both feature distribution shifts and sudden attack spikes automatically every 500 predictions.", BLUE),
        ("05", "End-to-end: CSV to live API",
         "Full pipeline: raw CSV → label mapping → SMOTE → 5 trained models → FastAPI → Docker → /drift/status in production. Ready to demo.", PURPLE),
        ("06", "Temporal split prevents data leakage",
         "Splitting by row order (time proxy) ensures the model is evaluated on 'future' traffic — not shuffled training data sneaking into the test set.", CYAN),
    ]
    cols = [0, 1, 2, 0, 1, 2]
    rows = [0, 0, 0, 1, 1, 1]
    for i, (num, title, body, col) in enumerate(takeaways):
        x = 0.4 + cols[i] * 3.1
        y = 1.05 + rows[i] * 2.12
        card_rect(s, x, y, 2.9, 1.98)
        accent_bar(s, x, y, 1.98, col)
        txb(s, num,   x + 0.18, y + 0.1,  0.55, 0.42, size=20, bold=True, color=col)
        txb(s, title, x + 0.18, y + 0.52, 2.6,  0.38, size=11, bold=True, color=WHITE)
        txb(s, body,  x + 0.18, y + 0.92, 2.6,  0.95, size=8.5, color=MUTED)


# ═════════════════════════════════════════════════════════════════════════════
#  MAIN
# ═════════════════════════════════════════════════════════════════════════════
def main():
    prs = new_prs()
    slide01(prs)
    slide02(prs)
    slide03(prs)
    slide04(prs)
    slide05(prs)
    slide06(prs)
    slide07(prs)
    slide08(prs)
    slide09(prs)
    slide10(prs)
    slide11(prs)
    slide12(prs)
    slide13(prs)
    slide14(prs)
    prs.save(OUT)
    print(f"DONE  Saved -> {OUT}")
    print(f"   14 slides  |  5 models  |  actual plots embedded")

if __name__ == "__main__":
    main()
